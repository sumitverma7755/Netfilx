from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9 (widescreen)
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define colors (Netflix-inspired)
NETFLIX_RED = RGBColor(229, 9, 20)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(20, 20, 20)

# Helper function to add a title slide
def add_title_slide(title, subtitle=None, background_color=BLACK):
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_color
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Format title
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = NETFLIX_RED
            run.font.size = Pt(60)
            run.font.bold = True
            run.font.name = 'Arial'
    
    # Add subtitle if provided
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Format subtitle
        for paragraph in subtitle_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = WHITE
                run.font.size = Pt(32)
                run.font.name = 'Arial'
    
    return slide

# Helper function to add a content slide
def add_content_slide(title, content, background_color=BLACK):
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_color
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Format title
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = NETFLIX_RED
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.name = 'Arial'
    
    # Add content
    content_shape = slide.placeholders[1]
    content_shape.text = content
    
    # Format content
    for paragraph in content_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.color.rgb = WHITE
            run.font.size = Pt(28)
            run.font.name = 'Arial'
    
    return slide

# Helper function to add a feature slide with bullet points
def add_feature_slide(title, features, background_color=BLACK):
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_color
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Format title
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = NETFLIX_RED
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.name = 'Arial'
    
    # Add bullet points
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    
    for i, feature in enumerate(features):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = feature
        p.level = 0
        
        # Format bullet points
        for run in p.runs:
            run.font.color.rgb = WHITE
            run.font.size = Pt(28)
            run.font.name = 'Arial'
    
    return slide

# Create slides for the Netfix app trailer

# 1. Intro slide
add_title_slide("NETFIX", "Your Ultimate Streaming Experience")

# 2. User Profiles slide
add_feature_slide("Personalized Profiles", [
    "✓ Multiple user profiles",
    "✓ Individual viewing preferences",
    "✓ Customized recommendations",
    "✓ Easy profile switching"
])

# 3. Movie Categories slide
add_feature_slide("Discover Content", [
    "✓ Browse by genre: Action, Comedy, Drama, Sci-Fi, Documentary",
    "✓ Trending now section",
    "✓ New releases",
    "✓ Award-winning titles"
])

# 4. Movie Details slide
add_feature_slide("Detailed Information", [
    "✓ Comprehensive movie details",
    "✓ Cast and crew information",
    "✓ User ratings and reviews",
    "✓ Similar recommendations"
])

# 5. Video Playback slide
add_feature_slide("Seamless Streaming", [
    "✓ High-definition playback",
    "✓ Adaptive streaming quality",
    "✓ Intuitive playback controls",
    "✓ Download for offline viewing"
])

# 6. Personalization slide
add_feature_slide("Smart Recommendations", [
    "✓ AI-powered content suggestions",
    "✓ Based on viewing history",
    "✓ Tailored to your preferences",
    "✓ Discover new favorites"
])

# 7. Closing slide
add_title_slide("NETFIX", "Stream Smarter. Download Now.")

# Save the presentation
output_path = os.path.join(os.getcwd(), "Netfix_App_Trailer.pptx")
prs.save(output_path)

print(f"Presentation created successfully at: {output_path}")
